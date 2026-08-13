"""Build the Meta Product Insights deck on top of the .potx template.

Strategy: the .potx is rewritten as a .pptx (Content Types swap), opened with
python-pptx, and we add slides using the template's own slide layouts so all
brand styling (fonts, colours, master backgrounds) is preserved.
"""

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
SRC_POTX = HERE / "template.potx"
INTERMEDIATE = HERE / "_intermediate.pptx"
OUTPUT = HERE / "Meta_Product_Insights_for_IT.pptx"
SHOTS = HERE / "screenshots"

TEMPLATE_CT = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
PRES_CT = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"

MONO_FONT = "Menlo"
CODE_PT = 10


def potx_to_pptx(src: Path, dst: Path) -> None:
    with zipfile.ZipFile(src, "r") as zin:
        files = {n: zin.read(n) for n in zin.namelist()}
    files["[Content_Types].xml"] = (
        files["[Content_Types].xml"].decode("utf-8").replace(TEMPLATE_CT, PRES_CT).encode("utf-8")
    )
    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n, d in files.items():
            zout.writestr(n, d)


def set_text(placeholder, text: str) -> None:
    tf = placeholder.text_frame
    tf.text = text


def set_bullets(placeholder, items: list[str]) -> None:
    tf = placeholder.text_frame
    tf.text = items[0]
    for line in items[1:]:
        p = tf.add_paragraph()
        p.text = line


def set_code(placeholder, code: str) -> None:
    """Render code as monospace, small, no bullet glyphs.

    Each line becomes its own paragraph so wrapping stays line-based.
    """
    tf = placeholder.text_frame
    tf.word_wrap = True
    lines = code.split("\n")
    tf.text = lines[0] if lines else ""
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
    for para in tf.paragraphs:
        # Kill bullet by zeroing indent; not all templates expose it, ignore on miss.
        para.level = 0
        for run in para.runs:
            run.font.name = MONO_FONT
            run.font.size = Pt(CODE_PT)


def ph_by_idx(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"placeholder idx={idx} not found in slide")


def insert_picture(slide, idx: int, image_path: Path) -> bool:
    """Drop a screenshot into the layout's picture placeholder.

    Returns True if the image was inserted, False if missing (caller can
    leave the empty placeholder visible as a TODO marker).
    """
    if not image_path.exists():
        return False
    ph = ph_by_idx(slide, idx)
    ph.insert_picture(str(image_path))
    return True


def remove_all_slides(prs: Presentation) -> None:
    """The .potx ships with sample slides — drop them before adding our own."""
    sldIdLst = prs.slides._sldIdLst  # CT_SlideIdList element
    for sldId in list(sldIdLst):
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


# ---------------------------------------------------------------------------
# Code snippets — kept as raw triple-strings so a reviewer reads exactly
# what ships in the repo. Trimmed to the load-bearing lines per topic.
# ---------------------------------------------------------------------------

CODE_SCHEDULER_TICK = '''\
# backend/app/services/scheduler.py
SCHEDULER_INTERVAL_S = 60

async def _scheduler_loop() -> None:
    while True:
        try:
            await _check_scheduled_jobs()
        except Exception as exc:
            logger.error("[Scheduler] loop error: %s", exc, exc_info=True)
        try:
            from .monitor import tick_due_monitors
            ran = await tick_due_monitors()
            if ran:
                logger.info("[Scheduler] Ran %d due monitor(s)", ran)
        except Exception as exc:
            logger.error("[Scheduler] Monitor tick failed: %s", exc, exc_info=True)
        await asyncio.sleep(SCHEDULER_INTERVAL_S)

async def _check_scheduled_jobs() -> None:
    now = datetime.utcnow()
    async with session_factory() as session:
        result = await session.execute(
            select(ScheduledJob).where(and_(
                ScheduledJob.enabled == True,
                ScheduledJob.nextRunAt <= now,
            ))
        )
        due_schedules = list(result.scalars().all())
    for schedule in due_schedules:
        await process_scheduled_job(schedule, trigger_type="auto")
'''

CODE_FB_ASYNC_POLL = '''\
# backend/app/facebook/insights.py
async def poll_report_status(report_run_id, access_token, on_progress=None):
    async with httpx.AsyncClient(timeout=30.0) as client:
        for _attempt in range(MAX_POLL_ATTEMPTS):
            url = f"https://graph.facebook.com/{GRAPH_API_VERSION}/{report_run_id}"
            resp = await client.get(url, params={"access_token": access_token})
            data = resp.json()
            if data.get("error"):
                raise RuntimeError(data["error"].get("message", "Poll error"))

            status  = data.get("async_status", "")
            percent = data.get("async_percent_completion", 0)
            if on_progress:
                await on_progress(percent)

            if status == "Job Completed":
                return {"success": True,
                        "async_report_url": data.get("async_report_url")}
            if status in ("Job Failed", "Job Skipped"):
                return {"success": False, "failure_reason": status}

            await asyncio.sleep(POLL_INTERVAL_MS / 1000)
    raise RuntimeError("Report generation timed out")
'''

CODE_CATALOG_BATCH_VERIFY = '''\
# backend/app/facebook/catalog.py — push
MAX_BATCH_SIZE = 3000
UPDATE_CONCURRENCY = 5

async def batch_update_products(catalog_id, requests, access_token, *, allow_upsert=False):
    batches   = [requests[i:i + MAX_BATCH_SIZE] for i in range(0, len(requests), MAX_BATCH_SIZE)]
    semaphore = asyncio.Semaphore(UPDATE_CONCURRENCY)

    async def _send(batch, idx):
        async with semaphore:
            if idx > 0: await asyncio.sleep(1.0)        # gentle pacing
            return await send_batch_request(catalog_id, batch, access_token,
                                            allow_upsert=allow_upsert, batch_index=idx)

    results = await asyncio.gather(*[_send(b, i) for i, b in enumerate(batches)],
                                   return_exceptions=True)
    # → aggregate handles + errors, return ParallelBatchResponse

# backend/app/services/report_worker.py — verify
async def verify_catalog_update(catalog_id, access_token, update_fields):
    base = f"https://graph.facebook.com/{API}/{catalog_id}/products"
    out  = {"total_catalog_products": 0, "fields": {}}
    async with httpx.AsyncClient(timeout=30) as c:
        total = (await c.get(base, params={"summary": "true", "limit": 0,
                                           "access_token": access_token})).json()
        out["total_catalog_products"] = total["summary"]["total_count"]
        for field, values in update_fields.items():
            f = {field: ({"eq": values[0]} if len(values) == 1 else {"is_any": values})}
            r = (await c.get(base, params={"summary": "true", "limit": 0,
                                           "filter": json.dumps(f),
                                           "access_token": access_token})).json()
            out["fields"][field] = {"matched": r["summary"]["total_count"],
                                    "expected_values": values}
    return out
'''

CODE_MONITOR_COUNT = '''\
# backend/app/facebook/product_set.py
async def fetch_product_set_count(product_set_id, access_token) -> int:
    """summary=true&limit=0 → one round-trip total count.
    68k-product set: ~6 min paginating → ~1 s here."""
    url    = f"{BASE_URL}/{product_set_id}/products"
    params = {"summary": "true", "limit": 0, "access_token": access_token}
    async with httpx.AsyncClient(timeout=30.0) as client:
        body = (await client.get(url, params=params)).json()
        return int(((body.get("summary") or {}).get("total_count")) or 0)

# backend/app/services/monitor.py
async def run_monitor(monitor_id: int, *, trigger_type: str = "auto") -> int:
    count = await fetch_product_set_count(product_set_id, access_token)
    prev  = (await s.execute(
        select(ProductSetSnapshot)
        .where(ProductSetSnapshot.monitorId == monitor_id,
               ProductSetSnapshot.status == "completed")
        .order_by(ProductSetSnapshot.takenAt.desc()).limit(1)
    )).scalar_one_or_none()
    snap = ProductSetSnapshot(monitorId=monitor_id, status="completed",
                              productCount=count, durationMs=duration_ms)
    # → delta = count - prev.productCount  → charted on Monitors page
'''


def main() -> None:
    potx_to_pptx(SRC_POTX, INTERMEDIATE)
    prs = Presentation(str(INTERMEDIATE))
    remove_all_slides(prs)

    # --- Slide 1: TITLE ---
    s = prs.slides.add_slide(prs.slide_layouts[0])
    set_text(ph_by_idx(s, 0), "Meta Product Insights")
    set_text(ph_by_idx(s, 1), "Catalog reporting & automation tool")
    set_text(ph_by_idx(s, 2), "For IT Review")
    set_text(ph_by_idx(s, 3), "William Lion — Product Marketing, Meta APAC")
    set_text(ph_by_idx(s, 4), "June 2026")

    # --- Slide 2: AGENDA ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_text(ph_by_idx(s, 0), "Agenda")
    set_bullets(
        ph_by_idx(s, 1),
        [
            "Why this tool exists",
            "Tech stack & architecture",
            "End-to-end data flow",
            "Core logic — code walkthrough",
            "Feature overview & UI",
            "Deployment options & ask from IT",
        ],
    )

    # --- Slide 3: MISSION ---
    s = prs.slides.add_slide(prs.slide_layouts[3])
    set_text(
        ph_by_idx(s, 0),
        "Turn 1M+ rows of FB product-level Insights into prioritised Catalog signals — on a schedule, with verification.",
    )
    set_text(
        ph_by_idx(s, 1),
        "Run weekly or daily, filter top performers by CTR/spend, push custom_label_* back to the Catalog, then verify how many products actually received the update.",
    )

    # --- Slide 4: TECH STACK ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_text(ph_by_idx(s, 0), "Built on a small, boring stack")
    set_text(ph_by_idx(s, 2), "Tech stack")
    set_bullets(
        ph_by_idx(s, 1),
        [
            "Frontend: React 19, Vite 7, TanStack Query, Tailwind, Radix UI, Recharts",
            "Backend: Python 3.12, FastAPI, SQLAlchemy 2 (async), APScheduler, httpx",
            "DB: SQLite (local) or TiDB Cloud (production) — same code, swap via DATABASE_URL",
            "External: Meta Marketing API v25.0 (Async Insights + Catalog Batch + Product Sets)",
            "Deploy: Single Docker image (Vite build → FastAPI static mount) → Cloud Run",
        ],
    )

    # --- Slide 5: ARCHITECTURE ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_text(ph_by_idx(s, 0), "System architecture")
    set_text(ph_by_idx(s, 2), "Architecture")
    set_bullets(
        ph_by_idx(s, 1),
        [
            "React SPA  —  served from FastAPI in production; Vite dev server locally",
            "FastAPI backend  —  REST API, in-process scheduler + job processor",
            "APScheduler  —  60-second tick, runs due cron schedules and monitors",
            "Meta Marketing API  —  Async Insights, Catalog Batch update, Product Set count",
            "Storage  —  SQLite file locally, TiDB Cloud (managed MySQL) in production",
            "No row-level data persisted  —  only DB summary (counts, spend, verify result)",
        ],
    )

    # --- Slide 6: DATA FLOW ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_text(ph_by_idx(s, 0), "How a scheduled run works")
    set_text(ph_by_idx(s, 2), "Data flow")
    set_bullets(
        ph_by_idx(s, 1),
        [
            "01  Trigger  —  APScheduler tick fires the cron at the configured time",
            "02  Fetch  —  POST FB Insights async report, poll until ready, download CSV (~290 MB)",
            "03  Filter  —  Parse CSV in-memory, apply CTR/spend filters, keep top performers",
            "04  Push  —  Catalog batch update (3000 items/req, 5 concurrent), retry on 5xx",
            "05  Verify  —  Catalog API summary query confirms matched count vs expected value",
            "06  Persist  —  Summary stats only — no raw payload written to disk or DB",
        ],
    )

    # ------------------------------------------------------------------
    # CORE LOGIC — code walkthrough
    # ------------------------------------------------------------------

    # --- Slide 7: section header for code ---
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_text(ph_by_idx(s, 0), "Core logic — code walkthrough")

    # --- Slide 8: Scheduler tick (60s loop) ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_text(ph_by_idx(s, 0), "Scheduler tick — 60-second loop")
    set_text(ph_by_idx(s, 2), "scheduler.py")
    set_code(ph_by_idx(s, 1), CODE_SCHEDULER_TICK)

    # --- Slide 9: FB async polling ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_text(ph_by_idx(s, 0), "FB async Insights polling")
    set_text(ph_by_idx(s, 2), "facebook/insights.py")
    set_code(ph_by_idx(s, 1), CODE_FB_ASYNC_POLL)

    # --- Slide 10: Catalog batch + verify ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_text(ph_by_idx(s, 0), "Catalog batch update + verification")
    set_text(ph_by_idx(s, 2), "catalog.py + report_worker.py")
    set_code(ph_by_idx(s, 1), CODE_CATALOG_BATCH_VERIFY)

    # --- Slide 11: Monitor count ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_text(ph_by_idx(s, 0), "Product Set Monitor — one-shot count + diff")
    set_text(ph_by_idx(s, 2), "product_set.py + monitor.py")
    set_code(ph_by_idx(s, 1), CODE_MONITOR_COUNT)

    # ------------------------------------------------------------------
    # PROOF + FEATURES
    # ------------------------------------------------------------------

    # --- Slide 12: PROOF — 3 stats ---
    s = prs.slides.add_slide(prs.slide_layouts[21])
    set_text(ph_by_idx(s, 0), "1.08M")
    set_text(ph_by_idx(s, 1), "Product-level rows fetched from FB Insights API")
    set_text(ph_by_idx(s, 4), "70,684")
    set_text(ph_by_idx(s, 5), "Products updated in Catalog after CTR ≥ 10% filter")
    set_text(ph_by_idx(s, 6), "3m 24s")
    set_text(ph_by_idx(s, 7), "End-to-end wall clock, errors = 0")
    set_text(ph_by_idx(s, 2), "Proof of execution")
    set_text(ph_by_idx(s, 3), "Real run, ad account 10153679704478431, 2026-06-01")

    # --- Slide 13: FEATURE OVERVIEW ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_text(ph_by_idx(s, 0), "Feature overview")
    set_bullets(
        ph_by_idx(s, 1),
        [
            "Schedules — Cron-based daily/weekly runs, multi-account support",
            "Report + Catalog — Pull Insights, filter, push custom_label_* in one job",
            "Product Set Monitor — Daily snapshot of product count, 30-day chart with delta",
            "Catalog Verification — Post-update API check confirms matched count",
            "Schedule History — Per-run timeline: status, duration, items, errors",
            "Token Vault — Per-user ads_management & catalog_management tokens; never logged",
        ],
    )

    # ------------------------------------------------------------------
    # UI SCREENSHOTS — drop PNGs into docs/screenshots/ and re-run.
    # Missing image → placeholder shows the layout's TODO frame.
    # ------------------------------------------------------------------

    # --- Slide 14: Schedules form ---
    s = prs.slides.add_slide(prs.slide_layouts[9])
    set_text(ph_by_idx(s, 0), "UI — Schedules form")
    set_text(ph_by_idx(s, 2), "/schedules")
    set_bullets(
        ph_by_idx(s, 1),
        [
            "Cron expression + timezone — preview of next run resolved server-side",
            "Per-user FB token picked from the Token Vault, never typed in the UI",
            "Filters (CTR ≥, spend ≥, custom label target) saved with the schedule",
            "Enable/disable toggle controls _check_scheduled_jobs eligibility",
        ],
    )
    insert_picture(s, 3, SHOTS / "schedules.png")

    # --- Slide 15: Run history — verify section ---
    s = prs.slides.add_slide(prs.slide_layouts[9])
    set_text(ph_by_idx(s, 0), "UI — Run history (catalog verification)")
    set_text(ph_by_idx(s, 2), "/reports")
    set_bullets(
        ph_by_idx(s, 1),
        [
            "Each run row expands to show total products vs matched per field",
            "matched value comes from verify_catalog_update → summary.total_count",
            "Gap between matched and expected is the headline reliability signal",
            "Duration, batch count, and error log live alongside the verify block",
        ],
    )
    insert_picture(s, 3, SHOTS / "reports.png")

    # --- Slide 16: Monitors — 30-day chart ---
    s = prs.slides.add_slide(prs.slide_layouts[9])
    set_text(ph_by_idx(s, 0), "UI — Product Set Monitors (30-day chart)")
    set_text(ph_by_idx(s, 2), "/monitors")
    set_bullets(
        ph_by_idx(s, 1),
        [
            "One snapshot per day via tick_due_monitors → fetch_product_set_count",
            "Recharts line plots productCount; delta vs prev snapshot annotated",
            "summary=true&limit=0 = one round-trip — 68k items in ~1 s, not ~6 min",
            "Failed snapshots stored with errorMessage; chart skips, keeps timeline",
        ],
    )
    insert_picture(s, 3, SHOTS / "monitors.png")

    # ------------------------------------------------------------------
    # DEPLOYMENT + ASK + CLOSING
    # ------------------------------------------------------------------

    # --- Slide 17: LOCAL vs CLOUD RUN ---
    s = prs.slides.add_slide(prs.slide_layouts[13])
    set_text(ph_by_idx(s, 0), "Local (dev / demo)")
    set_bullets(
        ph_by_idx(s, 1),
        [
            "SQLite file in backend/meta_insights.db",
            "uvicorn :8001 + vite :5174; reload on save",
            "Fires only while laptop awake & process running",
            "Zero infra cost; fastest iteration loop",
        ],
    )
    set_text(ph_by_idx(s, 3), "Cloud Run (production)")
    set_bullets(
        ph_by_idx(s, 4),
        [
            "TiDB Cloud as managed MySQL; schema auto-migrates",
            "Single container (Vite build → FastAPI static mount)",
            "Cloud Scheduler triggers /api/schedules/{id}/run on cron",
            "cpu-throttling + scale-to-zero ≈ USD $0–2 / month",
        ],
    )
    set_text(ph_by_idx(s, 2), "Deployment")
    try:
        set_text(ph_by_idx(s, 23), "Local dev vs Cloud Run")
    except KeyError:
        pass

    # --- Slide 18: ASK FROM IT ---
    s = prs.slides.add_slide(prs.slide_layouts[3])
    set_text(ph_by_idx(s, 0), "What we need from IT")
    set_text(
        ph_by_idx(s, 1),
        "1) GCP project + billing  ·  2) TiDB Cloud or Cloud SQL access  ·  "
        "3) FB long-lived token rotation policy  ·  4) Optional: domain + SSO",
    )

    # --- Slide 19: CLOSING ---
    s = prs.slides.add_slide(prs.slide_layouts[5])
    set_text(ph_by_idx(s, 0), "Thank you")

    prs.save(str(OUTPUT))
    print(f"OK — saved {OUTPUT} ({OUTPUT.stat().st_size:,} bytes), slides = {len(prs.slides)}")


if __name__ == "__main__":
    main()
